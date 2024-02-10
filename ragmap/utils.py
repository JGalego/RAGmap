"""
Helper functions to process documents,
project embeddings and plot the results.
"""

import os
import logging

from typing import (
    Any,
    Dict,
    List,
    Tuple,
    Union
)

import pandas as pd
import plotly.graph_objs as go

# File types
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

from .constants import (
    ModelProvider,
    DimensionReduction,
    plot_settings,
    plot3d_settings
)

logger = logging.getLogger(__name__)

def process_file(
    input: Any  # pylint: disable=redefined-builtin
    ) -> Tuple[List[str], Dict]:
    """
    Digests a file and returns text and metadata.
    """
    if isinstance(input, str) and os.path.isfile(input):
        _, f_ext = os.path.splitext(input)
        # PDF
        if f_ext == ".pdf":
            f = PdfReader(input)
        # DOCX
        elif f_ext == ".docx":
            f = Document(input)
        # PPTX
        elif f_ext == ".pptx":
            f = Presentation(input)
        else:
            raise NotImplementedError(f"{f_ext} files are not supported!")
    else:
        raise FileNotFoundError(f)

    # Extract text from file
    logger.info("Extracting text from file")
    if isinstance(f, PdfReader):
        texts = [p.extract_text().strip() for p in f.pages]
    elif isinstance(f, type(Document)):
        texts = [paragraph.text for paragraph in f.paragraphs]
    elif isinstance(f, type(Presentation)):
        texts = [shape.text for slide in f.slides \
                                for shape in slide.shapes \
                                    if hasattr(shape, "text")]

    # Retrieve document metadata
    # Adapted from https://stackoverflow.com/a/62022443
    logger.info("Retrieving document metadata")
    if isinstance(f, Union[type(Document), type(Presentation)]):
        metadata = {}
        fprops = f.core_properties
        metadata['author'] = fprops.author
        metadata['category'] = fprops.category
        metadata['comments'] = fprops.comments
        metadata['content_status'] = fprops.content_status
        metadata['created'] = fprops.created
        metadata['identifier'] = fprops.identifier
        metadata['keywords'] = fprops.keywords
        metadata['last_modified_by'] = fprops.last_modified_by
        metadata['language'] = fprops.language
        metadata['modified'] = fprops.modified
        metadata['subject'] = fprops.subject
        metadata['title'] = fprops.title
        metadata['version'] = fprops.version
    elif isinstance(f, PdfReader):
        metadata = f.metadata

    return texts, metadata

def plot_projections(
    df: pd.DataFrame,
    provider: ModelProvider,
    embedding_model: str,
    dimension_reduction: DimensionReduction,
    n_components: int) -> go.Figure:
    """
    Creates a plot of categorized embedding projections.
    """
    fig = go.Figure()
    for category in df['category'].unique():
        df_cat = df[df['category'] == category]
        # 2D
        if n_components == 2:
            category_settings = plot_settings[category]
            trace = go.Scatter(
                x=df_cat['x'],
                y=df_cat['y'],
                mode="markers",
                name=category.value,
                marker={
                    'color': category_settings['color'],
                    'opacity': category_settings['opacity'],
                    'symbol': category_settings['symbol'],
                    'size': category_settings['size'],
                    'line_width': 0
                },
                hoverinfo="text",
                text=df_cat['document_cleaned']
            )
        elif n_components == 3:
            category_settings = plot3d_settings[category]
            trace = go.Scatter3d(
                x=df_cat['x'],
                y=df_cat['y'],
                z=df_cat['z'],
                mode="markers",
                name=category.value,
                marker={
                    'color': category_settings['color'],
                    'opacity': category_settings['opacity'],
                    'symbol': category_settings['symbol'],
                    'size': category_settings['size'],
                    'line_width': 0
                },
                hoverinfo="text",
                text=df_cat['document_cleaned']
            )
        fig.add_trace(trace)

    fig.update_layout(
        title={
            'text': f"""
<sup>{provider.value} | {embedding_model} | {n_components}D {dimension_reduction.value.__name__}</sup>
""",
            'x': 0.5,
            'xanchor': 'center'
        },
        legend={
            'x': 0.5,
            'xanchor': "center",
            'yanchor': "bottom",
            'orientation': "h"
        }
    )
    return fig
