# pylint: disable=import-outside-toplevel
r"""
   ___  ___  _____
  / _ \/ _ |/ ___/_ _  ___ ____
 / , _/ __ / (_ /  ' \/ _ `/ _ \
/_/|_/_/ |_\___/_/_/_/\_,_/ .__/
                         /_/
            ___,
       _.-'` __|__
     .'  ,-:` \;',`'-,
    /  .'-;_,;  ':-;_,'.
   /  /;   '/    ,  _`.-\
  |  | '`. (`     /` ` \`|
  |  |:.  `\`-.   \_   / |
  |  |     (   `,  .`\ ;'|
   \  \     | .'     `-'/
    \  `.   ;/        .'
     '._ `'-._____.-'`
        `-.____|
          _____|_____
    jgs  /___________\

A simple Python module that helps visualize document chunks and queries in embedding space.
"""

import os
import logging
import uuid

from enum import Enum

from textwrap import wrap

from typing import (
    Any,
    Dict,
    List,
    Optional,
    Sequence,
    TypeVar,
    Union
)

import chromadb
import langchain
import numpy as np
import pandas as pd
import plotly.graph_objs as go
from pydantic import BaseModel

from chromadb.utils.embedding_functions import (
    AmazonBedrockEmbeddingFunction,
    OpenAIEmbeddingFunction,
    SentenceTransformerEmbeddingFunction
)

# File types
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# Dimensionality reduction
from openTSNE import TSNE
from sklearn.decomposition import PCA
from umap import UMAP

# Initialize logger
logger = logging.getLogger(__name__)

class ModelProvider(Enum):
    """
    Supported embedding model providers.
    """
    AMAZON_BEDROCK = "Amazon Bedrock ⛰️"
    HUGGING_FACE = "HuggingFace 🤗"
    OPENAI = "OpenAI"

class DimensionReduction(Enum):
    """
    Supported dimensionality reduction techniques.
    """
    UMAP = UMAP
    TSNE = TSNE
    PCA = PCA

class IndexCategory(Enum):
    """
    The different categories for each indexed chunk.
    """
    CHUNK = "Chunk"
    QUERY = "Query"
    RETRIEVED = "Retrieved"
    SUB_QUERY = "Sub-query"
    HYPOTHETICAL_ANSWER = "Hypothetical Answer"

# 2D
plot_settings = {
    IndexCategory.CHUNK: {
        'color': 'blue',
        'opacity': 0.5,
        'symbol': 'circle',
        'size': 10,
    },
    IndexCategory.QUERY: {
        'color': 'red',
        'opacity': 1,
        'symbol': 'cross',
        'size': 15,
    },
    IndexCategory.RETRIEVED: {
        'color': 'green',
        'opacity': 1,
        'symbol': 'star',
        'size': 15,
    },
    IndexCategory.SUB_QUERY: {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 15,
    },
    IndexCategory.HYPOTHETICAL_ANSWER: {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 15,
    },
}

# 3D
plot3d_settings = {
    IndexCategory.CHUNK: {
        'color': 'blue',
        'opacity': 0.5,
        'symbol': 'circle',
        'size': 10,
    },
   IndexCategory.QUERY: {
        'color': 'red',
        'opacity': 1,
        'symbol': 'circle',
        'size': 15,
    },
    IndexCategory.RETRIEVED: {
        'color': 'green',
        'opacity': 1,
        'symbol': 'diamond',
        'size': 10,
    },
    IndexCategory.SUB_QUERY: {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 10,
    },
    IndexCategory.HYPOTHETICAL_ANSWER: {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 10,
    },
}

class RAGmapError(Exception):
    """
    Raised any time RAGmap fails.
    """

# General
T = TypeVar("T")
OneOrMany = Union[T, List[T]]

# Vectors + Embeddings
Vector = Union[Sequence[float], Sequence[int]]
VectorDB = chromadb.Collection
Embedding = Vector
Embeddings = List[Embedding]

class RAGmap(BaseModel):
    """
    RAGmap class for visualizing and exploring embeddings
    """
    text_splitter: Any
    embedding_model: str
    provider: ModelProvider
    collection_metadata: Dict | None = None
    _embedding_function: Optional[chromadb.api.types.EmbeddingFunction] = None
    _vectordb: Optional[VectorDB] = None
    _last_id: int = 0

    def __init__(self, **data):
        logger.info("Initializing RAGmap")
        super().__init__(**data)
        self._init_embedding_function()
        self._init_vectordb()

    def _init_embedding_function(self) -> None:
        """
        Initializes the embedding function based on the model provider.
        """
        logger.info("Initializing embedding function")
        if self.provider == ModelProvider.AMAZON_BEDROCK:
            self._embedding_function = AmazonBedrockEmbeddingFunction(
                session=self.session,
                model_name=self.embedding_model,
            )
        elif self.provider == ModelProvider.HUGGING_FACE:
            self._embedding_function = SentenceTransformerEmbeddingFunction(
                model_name=self.embedding_model,
            )
        elif self.provider == ModelProvider.OPENAI:
            self._embedding_function = OpenAIEmbeddingFunction(
                api_key=os.environ['OPENAI_API_KEY'],
                model_name=self.embedding_model
            )

    def _init_vectordb(self) -> None:
        """
        Initializes the data store.
        """
        logger.info("Creating vector database")
        chroma_client = chromadb.Client()
        self._vectordb = chroma_client.create_collection(
            name=uuid.uuid4().hex,
            embedding_function=self._embedding_function,
            metadata=self.collection_metadata
        )

    def _split_text(
        self,
        text: str) -> List[str]:
        """
        Splits incoming text into smaller chunks.
        """
        if langchain.text_splitter.TextSplitter in type(self.text_splitter).__mro__:
            return self.text_splitter.split_text(text)
        if callable(self.text_splitter):
            return self.text_splitter(text)
        raise RAGmapError(f"Text splitter of type {type(self.text_splitter)} is not supported!")

    def load_file(
        self,
        f: Any) -> None:
        """
        Indexes and stores a single file in the vector DB.
        """
        if isinstance(f, str) and os.path.isfile(f):
            logger.info("Loading file %s", f)
            _, f_ext = os.path.splitext(f)
            # PDF
            if f_ext == ".pdf":
                f = PdfReader(f)
            # DOCX
            elif f_ext == ".docx":
                f = Document(f)
            # PPTX
            elif f_ext == ".pptx":
                f = Presentation(f)
            else:
                raise RAGmapError(f"{f_ext} files are not supported!")
        else:
            raise RAGmapError(f"File {f} not found!")

        # Extract text from file
        logger.info("Extracting text from file")
        if isinstance(f, PdfReader):
            texts = [p.extract_text().strip() for p in f.pages]
        elif isinstance(f, type(Document)):
            texts = [paragraph.text for paragraph in f.paragraphs]
        elif isinstance(f, type(Presentation)):
            texts = [shape.text for slide in f.slides \
                                    for shape in slide.shapes \
                                        if hasattr(shape, "text")]

        # Get document metadata
        # Adapted from https://stackoverflow.com/a/62022443
        logger.info("Retrieving document metadata")
        if isinstance(f, Union[type(Document), type(Presentation)]):
            metadata = {}
            fprops = f.core_properties
            metadata['author'] = fprops.author
            metadata['category'] = fprops.category
            metadata['comments'] = fprops.comments
            metadata['content_status'] = fprops.content_status
            metadata['created'] = fprops.created
            metadata['identifier'] = fprops.identifier
            metadata['keywords'] = fprops.keywords
            metadata['last_modified_by'] = fprops.last_modified_by
            metadata['language'] = fprops.language
            metadata['modified'] = fprops.modified
            metadata['subject'] = fprops.subject
            metadata['title'] = fprops.title
            metadata['version'] = fprops.version
        elif isinstance(f, PdfReader):
            metadata = f.metadata

        # Split text into chunks
        logger.info("Splitting text into chunks")
        text = "\n\n".join(texts)
        chunks = self._split_text(text)

        # Store everything in the vector database
        logger.info("Storing data in vector database")
        self._vectordb.add(
            ids=list(map(str, range(self._last_id, self._last_id + len(chunks)))),
            documents=chunks,
            metadatas=[metadata] * len(chunks)
        )
        self._last_id += len(chunks)

    def load_files(
        self,
        paths: OneOrMany[str]):
        """
        Indexes and stores one or multiple files in the vector DB.
        """
        for path in paths:
            self.load_file(path)

    def query(  # pylint: disable=too-many-arguments
        self,
        query_embeddings: OneOrMany[chromadb.api.types.Embedding] | \
                          OneOrMany[np.ndarray] | \
                          None = None,
        query_texts: OneOrMany[str] | None = None,
        n_results: int = 10,
        where: chromadb.api.types.Where | None = None,
        where_document: chromadb.api.types.WhereDocument | None = None):
        """
        Returns similar entries in the vector database for provided embeddings or texts.
        """
        return self._vectordb.query(
            query_texts=query_texts,
            query_embeddings=query_embeddings,
            n_results=n_results,
            where=where,
            where_document=where_document,
            include=[
                'documents',
                'embeddings',
                'metadatas'
            ]
        )

    def _fit_transform(
        self,
        dim_red_method: DimensionReduction = DimensionReduction.UMAP,
        n_components: int = 2,
        **kwargs):
        """
        Generates a dimensionality reduction transform.
        """
        logger.info("Retrieving vectors from datastore")
        embeddings = np.array(self._vectordb.get(include=['embeddings'])['embeddings'])

        # Fit embeddings
        logger.info("Fitting embeddings")
        dim_red = dim_red_method.value(
            n_components=n_components,
            **kwargs
        )
        transform = dim_red.fit(embeddings)
        return transform

    def _create_projections(
        self,
        dim_red_method: DimensionReduction = DimensionReduction.UMAP,
        n_components: int = 2,
        **kwargs):
        """
        Uses dimensionality reduction techniques to create embedding projections.
        """
        embeddings = np.array(self._vectordb.get(include=['embeddings'])['embeddings'])
        logger.info("Generating %iD %s projections", n_components, dim_red_method.value.__name__)
        return self._fit_transform(dim_red_method, n_components, **kwargs).transform(embeddings)

    def _create_figure(
        self,
        df: pd.DataFrame,
        dim_red_method: DimensionReduction,
        n_components: int):
        """
        Creates a Plotly figure from a Pandas dataframe.
        """
        fig = go.Figure()
        for category in df['category'].unique():
            df_cat = df[df['category'] == category]
            # 2D
            if n_components == 2:
                category_settings = plot_settings[category]
                trace = go.Scatter(
                    x=df_cat['x'],
                    y=df_cat['y'],
                    mode="markers",
                    name=category.value,
                    marker={
                        'color': category_settings['color'],
                        'opacity': category_settings['opacity'],
                        'symbol': category_settings['symbol'],
                        'size': category_settings['size'],
                        'line_width': 0
                    },
                    hoverinfo="text",
                    text=df_cat['document_cleaned']
                )
            elif n_components == 3:
                category_settings = plot3d_settings[category]
                trace = go.Scatter3d(
                    x=df_cat['x'],
                    y=df_cat['y'],
                    z=df_cat['z'],
                    mode="markers",
                    name=category.value,
                    marker={
                        'color': category_settings['color'],
                        'opacity': category_settings['opacity'],
                        'symbol': category_settings['symbol'],
                        'size': category_settings['size'],
                        'line_width': 0
                    },
                    hoverinfo="text",
                    text=df_cat['document_cleaned']
                )
            fig.add_trace(trace)

        fig.update_layout(
            title={
                'text': f"<sup>{self.provider.value} | {self.embedding_model} | {n_components}D {dim_red_method.value.__name__}</sup>",  # pylint: disable=line-too-long
                'x': 0.5,
                'xanchor': 'center'
            },
            legend={
                'x': 0.5,
                'xanchor': "center",
                'yanchor': "bottom",
                'orientation': "h"
            }
        )
        return fig

    def plot(
        self,
        dim_red_method: DimensionReduction = DimensionReduction.UMAP,
        n_components: int = 2,
        query: Dict | None = None,
        dim_red_method_config: Dict | None = None):
        """
        Creates a 2D or 3D plot of the reduced embedding space.
        """
        # Check projection components
        assert n_components in [2, 3], \
            f"Unsupported number of components (Got: {n_components}, Expected: 2 or 3)"

        # Retrieve data from vector store
        data = self._vectordb.get(
            include=['documents']
        )
        ids = data['ids']
        documents = data['documents']

        # Generate projections
        if dim_red_method_config is None:
            dim_red_method_config = {}
        projections = self._create_projections(
            dim_red_method,
            n_components=n_components,
            **dim_red_method_config
        )

        # Run the query and get the result IDs
        if query is not None:
            logger.info("Running query\n%s", query)
            retrieved_ids = []
            for result in self.query(**query)['ids']:
                retrieved_ids.extend(map(int, result))
        else:
            retrieved_ids = []

        # Prepare dataframe
        logger.info("Preparing dataframe")
        df = pd.DataFrame({
            'id': [int(id) for id in ids],
            'x': projections[:, 0],
            'y': projections[:, 1],
            'z': projections[:, 2] if n_components == 3 else None,
            'document_cleaned': [
                "<br>".join(wrap(doc, 100))
                    for doc in documents
            ],
            'category': IndexCategory.CHUNK
        })

        # Mark retrieved IDs
        df.loc[df['id'].isin(retrieved_ids), 'category'] = IndexCategory.RETRIEVED

        # Add query projections
        if query is not None:
            query_embeddings = query.get(
                'query_embeddings', self._embedding_function(
                        [query['query_texts']] \
                            if isinstance(query['query_texts'], str) \
                            else query['query_texts']
                )
            )
            query_projs = self._fit_transform(
                dim_red_method,
                n_components
            ).transform(query_embeddings)
            df_query = pd.DataFrame({
                'x': query_projs[:, 0],
                'y': query_projs[:, 1],
                'z': query_projs[:, 2] if n_components == 3 else None,
                'document_cleaned': query.get('query_texts', None),
                'category': IndexCategory.QUERY
            })
            df = pd.concat([df, df_query], axis=0)

        return self._create_figure(df, dim_red_method, n_components)
