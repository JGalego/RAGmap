r"""
   ___  ___  _____              
  / _ \/ _ |/ ___/_ _  ___ ____ 
 / , _/ __ / (_ /  ' \/ _ `/ _ \
/_/|_/_/ |_\___/_/_/_/\_,_/ .__/
                         /_/
            ___,
       _.-'` __|__
     .'  ,-:` \;',`'-,
    /  .'-;_,;  ':-;_,'.
   /  /;   '/    ,  _`.-\
  |  | '`. (`     /` ` \`|
  |  |:.  `\`-.   \_   / |
  |  |     (   `,  .`\ ;'|
   \  \     | .'     `-'/
    \  `.   ;/        .'
     '._ `'-._____.-'`
        `-.____|
          _____|_____
    jgs  /___________\

A simple Streamlit application that helps visualize document chunks and queries in embedding space.

Inspired by DeepLearning.ai's short course on 'Advanced Retrieval for AI with Chroma'
https://www.deeplearning.ai/short-courses/advanced-retrieval-for-ai/
and Gabriel Chua's award-winning RAGxplorer
https://github.com/gabrielchua/RAGxplorer/
"""

import io
import json
import os
import uuid

from textwrap import wrap

import boto3
import botocore
import chromadb

import numpy as np
import pandas as pd
import plotly.graph_objs as go
import streamlit as st

from chromadb.utils.embedding_functions import (
    AmazonBedrockEmbeddingFunction,
    SentenceTransformerEmbeddingFunction
)

from huggingface_hub.utils import RepositoryNotFoundError

from langchain.prompts import ChatPromptTemplate

from langchain.text_splitter import (
    RecursiveCharacterTextSplitter
)

from langchain_community.chat_models import BedrockChat

from langchain_core.output_parsers import StrOutputParser

from openTSNE import TSNE

from sklearn.decomposition import PCA

from umap import UMAP

#############
# Constants #
#############

# State
state_vars = [
    "collection",
    "documents",
    "projections",
    "query_projections",
    "retrieved_ids"
]

# 2D
plot_settings = {
    'chunk': {
        'color': 'blue',
        'opacity': 0.5,
        'symbol': 'circle',
        'size': 10,
    },
    'query': {
        'color': 'red',
        'opacity': 1,
        'symbol': 'cross',
        'size': 15,
    },
    'retrieved': {
        'color': 'green',
        'opacity': 1,
        'symbol': 'star',
        'size': 15,
    },
    'sub-query': {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 15,
    },
    'hypothetical answer': {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 15,
    },
}

# 3D
plot3d_settings = {
    'chunk': {
        'color': 'blue',
        'opacity': 0.5,
        'symbol': 'circle',
        'size': 10,
    },
    'query': {
        'color': 'red',
        'opacity': 1,
        'symbol': 'circle',
        'size': 15,
    },
    'retrieved': {
        'color': 'green',
        'opacity': 1,
        'symbol': 'diamond',
        'size': 10,
    },
    'sub-query': {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 10,
    },
    'hypothetical answer': {
        'color': 'purple',
        'opacity': 1,
        'symbol': 'square',
        'size': 10,
    },
}

################
# Helper Utils #
################

st.cache_data()
def list_embedding_models():
    """
    Returns a list of embedding models available in Amazon Bedrock
    """
    try:
        # Filter out provisioned throughput only models
        # https://docs.aws.amazon.com/bedrock/latest/userguide/prov-throughput-models.html
        return [model for model in bedrock.list_foundation_models(
                                        byOutputModality='EMBEDDING')['modelSummaries'] \
                                            if model['inferenceTypesSupported'] != ["PROVISIONED"]]
    except botocore.exceptions.ClientError as error:
        st.error(error)
        return []


st.cache_data()
def list_text_models():
    """
    Returns a list of text models available in Amazon Bedrock
    """
    try:
        return bedrock.list_foundation_models(
            byOutputModality='TEXT')['modelSummaries']
    except botocore.exceptions.ClientError as error:
        st.error(error)
        return []


st.cache_data()
def model2table(model):
    """
    Turns a model summary into a table
    """
    if model is None:
        return None

    model_status = model['modelLifecycle']['status']
    model_status = "✅" if model_status == "ACTIVE" else "👴" if model_status == "LEGACY" else "❓"

    input_modalities = model['inputModalities']
    input_modalities = "".join(
        map(
            lambda inp: "💬" if inp == "TEXT" else "🖼️" if inp == "IMAGE" else "❓",
            input_modalities
        )
    )

    return f"""
<table style="margin: 0px auto;">
    <tr>
        <td>
            <b>Model ID</b>
        </td>
        <td>
            <tt>{model['modelId']}</tt>
        </td>
    </tr>
    <tr>
        <td>
            <b>Model Name</b>
        </td>
        <td>
            {model['modelName']}
        </td>
    </tr>
    <tr>
        <td>
            <b>Input Modalities</b>
        </td>
        <td>
            {input_modalities}
        </td>
    </tr>
    <tr>
        <td>
            <b>Output Modalities</b>
        </td>
        <td>
            {", ".join(model['outputModalities'])}
        </td>
    </tr>
    <tr>
        <td>
            <b>Inference Types</b>
        </td>
        <td>
            {", ".join(model['inferenceTypesSupported'])}
        </td>
    </tr>
    <tr>
        <td>
            <b>Model Lifecycle</b>
        </td>
        <td>
            {model_status}
        </td>
    </tr>
</table>

#####
"""


def process_document(filename):
    """
    Loads and extracts text from a document
    """
    # PDF
    mime_type = filename.type
    if mime_type == "application/pdf":
        from PyPDF2 import PdfReader  # pylint: disable=import-outside-toplevel
        doc = PdfReader(filename)
        texts = [p.extract_text().strip() for p in doc.pages]
    # DOCX
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        from docx import Document  # pylint: disable=import-outside-toplevel
        doc = Document(filename)
        texts = [paragraph.text for paragraph in doc.paragraphs]
    # PPTX
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation  # pylint: disable=import-outside-toplevel
        doc = Presentation(filename)
        texts = [shape.text for slide in doc.slides \
                                for shape in slide.shapes \
                                    if hasattr(shape, "text")]

    # Filter out empty strings
    texts = [text for text in texts if text]
    return texts


def chunk_texts(texts):
    """
    Breaks down input texts into chunks
    """
    character_splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", ". ", " ", ""],
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    chunks = character_splitter.split_text('\n\n'.join(texts))
    return chunks


def embedding_function():
    """
    Specifies the embedding function
    """
    if model_provider == "Amazon Bedrock ⛰️":
        st.session_state.embedding_function = AmazonBedrockEmbeddingFunction(
            session=session,
            model_name=embedding_model['modelId']
        )
    elif model_provider == "HuggingFace 🤗":
        try:
            st.session_state.embedding_function = SentenceTransformerEmbeddingFunction(
                model_name=embedding_model
            )
        except RepositoryNotFoundError as error:
            st.error(error)


def build_collection():
    """
    Creates and populates a Chroma collection
    """
    # Load and extract text from PDF document
    with st.spinner("Loading document"):
        texts = process_document(uploaded_file)

    # Split text into chunks
    with st.spinner("Splitting text into chunks"):
        chunks = chunk_texts(texts)

    # Create and populate chroma collection
    with st.spinner("Creating collection"):
        if os.path.exists("embeddings"):
            # Initialize a persistent chroma client
            # https://docs.trychroma.com/usage-guide#initiating-a-persistent-chroma-client
            chroma_client = chromadb.PersistentClient("embeddings")
        else:
            chroma_client = chromadb.Client()
        document_name = uuid.uuid4().hex
        collection = chroma_client.create_collection(
            document_name,
            embedding_function=st.session_state.embedding_function
        )
    with st.spinner("Populating collection"):
        collection.add(
            ids=list(map(str, range(len(chunks)))),
            documents=chunks
        )

    st.session_state.collection = collection


def get_embeddings(text):
    """
    Transforms input text into embeddings
    """
    text_embeddings = st.session_state.embedding_function([text])
    return text_embeddings


def project_embeddings(embeddings, transform):
    """
    Projects text embeddings using a transform
    """
    if isinstance(embeddings, list):
        embeddings = np.array(embeddings)
    projs = np.empty((len(embeddings), n_components))
    projs = transform.transform(embeddings)
    return projs


def create_projections():
    """
    Transforms document embeddings into projections
    """
    if st.session_state.collection is None:
        return

    # Get documents and embeddings
    with st.spinner("Retrieving document and embeddings"):
        res = st.session_state.collection.get(
            include=['documents', 'embeddings']
        )
        st.session_state.ids = res['ids']
        st.session_state.embeddings = np.array(res['embeddings'])
        st.session_state.documents = res['documents']

    # Fit projection transform
    with st.spinner("Fitting embeddings"):
        if st.session_state.dim_redux == "UMAP":
            transform = UMAP(
                random_state=0, transform_seed=0, n_components=st.session_state.n_components
            ).fit(st.session_state.embeddings)
        elif st.session_state.dim_redux == "t-SNE":
            transform = TSNE(
                random_state=0, n_components=st.session_state.n_components
            ).fit(st.session_state.embeddings)
        elif st.session_state.dim_redux == "PCA":
            transform = PCA(
                random_state=0, n_components=st.session_state.n_components
            ).fit(st.session_state.embeddings)
        st.session_state.transform = transform

    # Get projections
    with st.spinner("Generating projections"):
        st.session_state.projections = transform.transform(st.session_state.embeddings)


def plot_projections(df_projs):
    """
    Generates a 2D or 3D plot of embedding projections
    """
    fig = go.Figure()
    for category in df['category'].unique():
        df_cat = df_projs[df_projs['category'] == category]
        # 2D
        if n_components == 2:
            category_settings = plot_settings[category]
            trace = go.Scatter(
                x=df_cat['x'],
                y=df_cat['y'],
                mode='markers',
                name=category,
                marker={
                    'color': category_settings['color'],
                    'opacity': category_settings['opacity'],
                    'symbol': category_settings['symbol'],
                    'size': category_settings['size'],
                    'line_width': 0
				},
                hoverinfo='text',
                text=df_cat['document_cleaned']
            )
        # 3D
        elif n_components == 3:
            category_settings = plot3d_settings[category]
            trace = go.Scatter3d(
                x=df_cat['x'],
                y=df_cat['y'],
                z=df_cat['z'],
                mode="markers",
                name=category,
                marker={
                    'color': category_settings['color'],
                    'opacity': category_settings['opacity'],
                    'symbol': category_settings['symbol'],
                    'size': category_settings['size'],
                    'line_width': 0
				},
                hoverinfo='text',
                text=df_cat['document_cleaned']
            )

        fig.add_trace(trace)

    fig.update_layout(
        title={
            'text': f"{uploaded_file.name} <br><sup>{model_provider} | {embedding_model['modelName'] if model_provider == 'Amazon Bedrock ⛰️' else embedding_model if model_provider == 'HuggingFace 🤗' else None} | ({chunk_size}, {chunk_overlap}) chunks | {n_components}D {dim_redux}</sup>",  # pylint: disable=line-too-long
            'x': 0.5,
            'xanchor': 'center'
        },
        legend={
			'x': 0.5,
            'xanchor': "center",
            'yanchor': "bottom",
            'orientation': "h"
		}
    )

    st.plotly_chart(fig, use_container_width=True)
    return fig


def multiple_queries_expansion(user_query, df_query_original, model_id='anthropic.claude-v2:1'):
    """
    Expands a user query by creating multiple sub-queries
    """
    prompt = f"""
System: Given a query, your task is to generate 3 to 5 simple sub-queries related to the original query. These sub-queries must be short. Format your reply in JSON with numbered keys. Skip the preamble.

Human: "{user_query}"

Assistant:
"""  # pylint: disable=line-too-long
    body = body = json.dumps({
        "prompt": prompt, "max_tokens_to_sample": 1000
    })
    try:
        response = bedrock_runtime.invoke_model(
            modelId=model_id,
            body=body,
            accept="application/json",
            contentType="application/json"
        )
        output = json.loads(response.get("body").read())['completion']
        st.session_state.query_expansions = list(json.loads(output).values())
        st.session_state.query_expansion_projections = np.array([
            project_embeddings(
                get_embeddings(expansion),
                st.session_state.transform
            ) for expansion in st.session_state.query_expansions
        ])

        df_query_expansions = pd.DataFrame({
            'x': st.session_state.query_expansion_projections[:, 0, 0],
            'y': st.session_state.query_expansion_projections[:, 0, 1],
            'z': st.session_state.query_expansion_projections[:, 0, 2] \
                    if n_components == 3 else None,
            'document_cleaned': st.session_state.query_expansions,
            'category': ["sub-query"] * len(st.session_state.query_expansions),
        })

        df_query_expanded = pd.concat([df_query_original, df_query_expansions])
        return st.session_state.query_expansions, df_query_expanded
    except botocore.exceptions.ClientError as error:
        st.error(error)
        st.warning("Failed to expand query, falling back to the naive approach!")
        return [user_query], df_query_original
    except json.decoder.JSONDecodeError:
        st.warning("Failed to expand query, falling back to the naive approach!")
        return [user_query], df_query_original


def generated_answer_expansion(user_query, df_query_original, model_id='anthropic.claude-v2:1'):
    """
    Expands a user query by generating an hypothetical answer

    References:
    + (Gao et al., 2022) Precise Zero-Shot Dense Retrieval without Relevance Labels
    https://paperswithcode.com/paper/precise-zero-shot-dense-retrieval-without
    """
    system_msg = "Given a query, your task is to generate a template for the hypothetical answer. Do not include any facts, and instead label them as <PLACEHOLDER>. Skip the preamble."  # pylint: disable=line-too-long
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_msg),
        ("human", "\"{user_query}\"")
    ])
    model = BedrockChat(model_id=model_id)
    output_parser = StrOutputParser()
    chain = prompt | model | output_parser
    try:
        output = chain.invoke({"user_query": user_query})
        st.session_state.query_expansions = [output]
        st.session_state.query_expansion_projections = np.array([
            project_embeddings(
                get_embeddings(expansion),
                st.session_state.transform
            ) for expansion in st.session_state.query_expansions
        ])

        df_query_expansions = pd.DataFrame({
            'x': st.session_state.query_expansion_projections[:, 0, 0],
            'y': st.session_state.query_expansion_projections[:, 0, 1],
            'z': st.session_state.query_expansion_projections[:, 0, 2] \
                    if n_components == 3 else None,
            'document_cleaned': [
                "<br>".join(wrap(doc, 100)) for doc in st.session_state.query_expansions
            ],
            'category': ["hypothetical answer"] * len(st.session_state.query_expansions),
        })

        df_query_expanded = pd.concat([df_query_original, df_query_expansions])
        return st.session_state.query_expansions, df_query_expanded
    except botocore.exceptions.ClientError as error:
        st.error(error)
    except json.decoder.JSONDecodeError:
        pass
    st.warning("Model failed to expand query, falling back to the naive approach!")
    return [user_query], df_query_original


def initialize():
    """
    Pre-initializes the session state
    """
    for key in state_vars:
        if key not in st.session_state:
            st.session_state[key] = None


def reset():
    """
    Resets session state and re-runs the application
    """
    for key in state_vars:
        st.session_state[key] = None

########
# Main #
########

st.title("RAGMap 🗺️🔍")
st.text("From meaning to vectors and back...")

# Initialize session state
initialize()

# Initialize session
session = boto3.Session()

# Initialize bedrock clients
bedrock = session.client('bedrock')
bedrock_runtime = session.client('bedrock-runtime')

st.markdown("### 1. Upload a document 📄")

uploaded_file = st.file_uploader(
    label="Upload a file",
    label_visibility="hidden",
    type=["pdf", "docx", "pptx"],
)

st.markdown("### 2. Build a vector database 💫")

chunk_size = st.number_input(
    label="Chunk Size",
    min_value=1,
    max_value=2000,
    value=256,
    step=1,
    key="chunk_size",
    help="Number of tokens in each chunk",
)

chunk_overlap = st.number_input(
    label="Chunk Overlap",
    min_value=0,
    max_value=1000,
    value=0,
    step=1,
    key="chunk_overlap",
    help="Number of tokens shared between consecutive chunks to maintain context",
)

model_provider = st.radio(
    label="Model Provider",
    options=[
        "Amazon Bedrock ⛰️",
        "HuggingFace 🤗"
    ],
    horizontal=True,
    index=0,
    on_change=reset,
    key="model_provider",
    help="Where the model comes from",
)

if model_provider == "Amazon Bedrock ⛰️":
    embedding_model = st.selectbox(
        label="Embedding Model",
        options=list_embedding_models(),
        index=0,
        format_func=lambda option: f"{option['modelName']} ({option['modelId']})",
        key="embedding_model",
        help="The model used to encapsulate information \
            into dense representations in a multi-dimensional space",
    )
    st.markdown(model2table(embedding_model), unsafe_allow_html=True)
elif model_provider == "HuggingFace 🤗":
    embedding_model = st.text_input(
        label="Embedding Model",
        key="embedding_model",
        placeholder="Enter the model name e.g. all-MiniLM-L6-v2",
        help="The model used to encapsulate information \
            into dense representations in a multi-dimensional space",
    )

dim_redux = st.radio(
    label="Dimensionality Reduction",
    options=[
        "UMAP",
        "t-SNE",
        "PCA"
    ],
    captions=[
        "Uniform Manifold Approximation and Projection",
        "t-distributed Stochastic Neighbor Embedding",
        "Principal Component Analysis"
    ],
    horizontal=True,
    index=0,
    on_change=create_projections,
    key="dim_redux",
    help="The algorithm or technique used for dimensionality reduction",
)

# For more information, see
# https://umap-learn.readthedocs.io/en/latest/parameters.html#n-components
# https://opentsne.readthedocs.io/en/stable/api/index.html#openTSNE.TSNE
n_components = st.radio(
    label="Projection Components",
    options=[2, 3],
    horizontal=True,
    index=0,
    format_func=lambda option: f"{option}D",
    on_change=create_projections,
    key="n_components",
    help="The dimensionality of the reduced embedding space",
)

if embedding_model is not None:
    embedding_function()

if st.button(label="Build"):
    if uploaded_file is None:
        st.error("No file uploaded!")
    elif model_provider == "HuggingFace 🤗" and \
       len(embedding_model) == 0:
        st.error("Model name must not be empty!")
    elif embedding_model is None:
        st.error("No model selected!")
    else:
        if st.session_state.collection is None:
            build_collection()

        if st.session_state.projections is None:
            create_projections()

if st.session_state.projections is not None:

    st.markdown("### 3. Explore the embedding space 🖼️")

    df = pd.DataFrame({
        'id': [int(id) for id in st.session_state.ids],
        'x': st.session_state.projections[:, 0],
        'y': st.session_state.projections[:, 1],
        'z': st.session_state.projections[:, 2] if n_components == 3 else None,
        'document_cleaned': [
            "<br>".join(wrap(doc, 100))
                for doc in st.session_state.documents
        ],
        'category': "chunk",
    })

    n_results = st.number_input(
        label="Number of results",
        min_value=1,
        max_value=10,
        value=5,
        key='n_results',
    )

    retrieval_strategy = st.radio(
        label="Retrieval Strategy",
        options=[
            "Naive",
            "Multiple Queries",
            "Generated Answer",
        ],
        captions=[
            "Retrieves chunks with high similarity to the user query",
            "Expands the original query by generating additional sub-queries",
            "Expands the original query by generating an hypothetical answer",
        ],
        key='retrieval_strategy',
    )

    query = st.text_input(
        label="Query",
        key='query',
    )

    if st.button(label="Search"):
        # Get query projections
        st.session_state.query_projections = project_embeddings(
            get_embeddings(query),
            st.session_state.transform
        )
        df_query = pd.DataFrame({
            'x': st.session_state.query_projections[:, 0],
            'y': st.session_state.query_projections[:, 1],
            'z': st.session_state.query_projections[:, 2] if n_components == 3 else None,
            'document_cleaned': query,
            'category': "query",
        })

        # Use the query as is or expand it
        if retrieval_strategy == "Naive":
            chroma_query = query
        elif retrieval_strategy == "Multiple Queries":
            chroma_query, df_query = multiple_queries_expansion(query, df_query)
        elif retrieval_strategy == "Generated Answer":
            chroma_query, df_query = generated_answer_expansion(query, df_query)

        # Search query and process results
        results = st.session_state.collection.query(
            query_texts=chroma_query,
            n_results=n_results,
            include=['documents', 'embeddings']
        )

        st.session_state.retrieved_ids = []
        for ids in results['ids']:
            st.session_state.retrieved_ids.extend(ids)
        st.session_state.retrieved_ids = map(int, st.session_state.retrieved_ids)
        df.loc[df['id'].isin(st.session_state.retrieved_ids), 'category'] = "retrieved"

        # Append query projections
        df = pd.concat([df, df_query], axis=0)

    # Display all projections
    st.markdown("#### Projections")
    proj_plot = plot_projections(df)
    buffer = io.StringIO()
    proj_plot.write_html(buffer, include_plotlyjs="cdn")
    html_bytes = buffer.getvalue().encode()

    # Download plot as HTML
    st.download_button(
        label="Download HTML",
        data=html_bytes,
        file_name=f"{uploaded_file.name.split('.')[0]}.html",
        mime="text/html",
    )

    if st.session_state.retrieved_ids is not None:
        # Display query results
        st.markdown("#### Results")
        st.dataframe(
            df[df['category'] == "retrieved"].drop(['category'], axis=1),
            use_container_width=True,
            hide_index=True,
            column_config={
                'x': st.column_config.NumberColumn(
                    width="small"
                ),
                'document_cleaned': st.column_config.TextColumn(
                    "chunk",
                    width="large"
                )
            },
        )

        # Download query results as CSV
        st.download_button(
            label="Download CSV",
            data=df.to_csv(),
            file_name="results.csv",
            mime="text/csv",
        )
